"""Unit tests for Deck Refresh (no Cursor SDK / LibreOffice required)."""
from __future__ import annotations

import json
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from app.refresh.advise import advice_from_changelog
from app.refresh.compliance import scan_pptx
from app.refresh.keywords import find_keyword_hits
from app.refresh.refresh_deck import refresh_presentation


ROOT = Path(__file__).resolve().parents[3]
CONFIG = ROOT / "config"


def _make_pptx(path: Path, *, text: str = "Hello ST", fill=None, font_name: str = "Comic Sans MS") -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(1))
    tf = shape.text_frame
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.text = text
    run = tf.paragraphs[0].runs[0]
    run.font.name = font_name
    run.font.size = Pt(40)
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    if fill is not None:
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2), Inches(4), Inches(1.2)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.text_frame.text = "Card"
        box.text_frame.paragraphs[0].runs[0].font.name = font_name
    prs.save(str(path))


def test_scan_rejects_non_zip(tmp_path: Path):
    bad = tmp_path / "x.pptx"
    bad.write_bytes(b"not-a-zip")
    result = scan_pptx(bad, CONFIG)
    assert result.blocked
    assert result.reason == "protected_file"


def test_scan_accepts_pptx(tmp_path: Path):
    pptx = tmp_path / "ok.pptx"
    _make_pptx(pptx)
    result = scan_pptx(pptx, CONFIG)
    assert not result.blocked
    assert any(layer.id == "file_integrity" and layer.status == "pass" for layer in result.layers)


def test_refresh_does_not_change_text(tmp_path: Path):
    src = tmp_path / "in.pptx"
    dest = tmp_path / "out.pptx"
    _make_pptx(src, text="Keep this exact text 123")
    log = refresh_presentation(src, dest, tmp_path / "changelog.json")
    assert dest.is_file()
    prs = Presentation(str(dest))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    assert "Keep this exact text 123" in texts
    assert log.slide_count == 1
    assert (tmp_path / "changelog.json").is_file()


def test_refresh_sets_arial_and_disables_autofit(tmp_path: Path):
    src = tmp_path / "in.pptx"
    dest = tmp_path / "out.pptx"
    _make_pptx(src)
    refresh_presentation(src, dest, tmp_path / "cl.json")
    prs = Presentation(str(dest))
    slide = prs.slides[0]
    shape = slide.shapes[0]
    assert shape.text_frame.auto_size == MSO_AUTO_SIZE.NONE
    names = {run.font.name for para in shape.text_frame.paragraphs for run in para.runs}
    assert "Arial" in names


def test_refresh_preserves_page_count(tmp_path: Path):
    src = tmp_path / "multi.pptx"
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for i in range(3):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        tb.text_frame.text = f"Slide {i+1}"
    prs.save(str(src))
    dest = tmp_path / "out.pptx"
    log = refresh_presentation(src, dest)
    assert log.slide_count == 3
    assert len(Presentation(str(dest)).slides) == 3


def test_advice_from_changelog_zh():
    md = advice_from_changelog(
        {
            "slide_count": 2,
            "summary": {"color": 3, "align": 1},
            "skips": [{"slide": 1, "shape_id": 1, "reason": "chart_skipped"}],
            "warnings": [],
            "changes": [{"slide": 1, "shape_id": 1, "kind": "color", "detail": "x"}],
        },
        lang="zh",
    )
    assert "已自动修正" in md
    assert "建议您考虑" in md


def test_keyword_hit_when_enabled(tmp_path: Path, monkeypatch):
    # Temporarily point at a config with keywords enabled via rewriting settings load
    pptx = tmp_path / "sec.pptx"
    _make_pptx(pptx, text="CONFIDENTIAL draft for review")
    cfg_dir = tmp_path / "cfg"
    cfg_dir.mkdir()
    (cfg_dir / "compliance.json").write_text(
        json.dumps(
            {
                "purview_enabled": False,
                "keyword_enabled": True,
                "block_on_keyword_hit": True,
                "block_on_protected_file": True,
            }
        ),
        encoding="utf-8",
    )
    (cfg_dir / "compliance_keywords.json").write_text(
        (CONFIG / "compliance_keywords.json").read_text(encoding="utf-8"),
        encoding="utf-8",
    )
    hits = find_keyword_hits(pptx, cfg_dir)
    assert any(h["phrase"].upper() == "CONFIDENTIAL" for h in hits)
    result = scan_pptx(pptx, cfg_dir)
    assert result.require_ack
    assert result.reason == "keyword_hit"
