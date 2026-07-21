"""Keyword compliance scan for Deck Refresh (config-driven)."""
from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from pptx import Presentation


def load_keyword_config(config_dir: Path) -> dict[str, Any]:
    path = config_dir / "compliance_keywords.json"
    if not path.is_file():
        return {}
    return json.loads(path.read_text(encoding="utf-8"))


def _collect_text(prs: Presentation, scan_core: bool) -> str:
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                parts.append(shape.text_frame.text or "")
            if getattr(shape, "has_table", False):
                try:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            parts.append(cell.text or "")
                except Exception:  # noqa: BLE001
                    pass
    if scan_core:
        try:
            cp = prs.core_properties
            for attr in ("title", "subject", "keywords", "comments", "category"):
                val = getattr(cp, attr, None)
                if val:
                    parts.append(str(val))
        except Exception:  # noqa: BLE001
            pass
    return "\n".join(parts)


def find_keyword_hits(path: Path, config_dir: Path) -> list[dict[str, str]]:
    cfg = load_keyword_config(config_dir)
    phrases = list(cfg.get("phrases_en") or []) + list(cfg.get("phrases_zh") or [])
    if not phrases:
        return []
    match = cfg.get("match") or {}
    case_insensitive = bool(match.get("case_insensitive", True))
    # Longest first to prefer longer phrases in reporting
    if match.get("prefer_longest_phrase_first", True):
        phrases = sorted(phrases, key=len, reverse=True)

    prs = Presentation(str(path))
    text = _collect_text(prs, bool(cfg.get("scan_core_properties", True)))
    hay = text.lower() if case_insensitive else text
    hits: list[dict[str, str]] = []
    seen: set[str] = set()
    for phrase in phrases:
        needle = phrase.lower() if case_insensitive else phrase
        if not needle or needle in seen:
            continue
        if needle in hay:
            seen.add(needle)
            hits.append({"phrase": phrase})
    return hits
