"""Deterministic Deck Refresh rule engine (standard tier).

Applies visual-only changes via python-pptx. Never edits text content.
Writes refresh_changelog.json as the source of truth for changes.
"""
from __future__ import annotations

import json
import math
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from . import rules_config as rc


@dataclass
class ChangeEntry:
    slide: int
    shape_id: int
    kind: str
    detail: str


@dataclass
class SkipEntry:
    slide: int
    shape_id: int | None
    reason: str


@dataclass
class Changelog:
    version: int = 1
    tier: str = "standard"
    source: str = ""
    slide_count: int = 0
    changes: list[ChangeEntry] = field(default_factory=list)
    skips: list[SkipEntry] = field(default_factory=list)
    summary: dict[str, int] = field(default_factory=dict)
    warnings: list[str] = field(default_factory=list)

    def add(self, slide: int, shape_id: int, kind: str, detail: str) -> None:
        self.changes.append(ChangeEntry(slide, shape_id, kind, detail))

    def skip(self, slide: int, shape_id: int | None, reason: str) -> None:
        self.skips.append(SkipEntry(slide, shape_id, reason))

    def finalize(self) -> None:
        counts: dict[str, int] = {}
        for c in self.changes:
            counts[c.kind] = counts.get(c.kind, 0) + 1
        self.summary = counts

    def to_dict(self) -> dict[str, Any]:
        self.finalize()
        return {
            "version": self.version,
            "tier": self.tier,
            "source": self.source,
            "slide_count": self.slide_count,
            "summary": self.summary,
            "warnings": self.warnings,
            "changes": [asdict(c) for c in self.changes],
            "skips": [asdict(s) for s in self.skips],
        }


def _rgb_tuple(color: RGBColor | None) -> tuple[int, int, int] | None:
    if color is None:
        return None
    try:
        return (int(color[0]), int(color[1]), int(color[2]))
    except Exception:  # noqa: BLE001
        return None


def _dist(a: tuple[int, int, int], b: tuple[int, int, int]) -> float:
    return math.sqrt(sum((x - y) ** 2 for x, y in zip(a, b)))


def _nearest_fill(rgb: tuple[int, int, int]) -> RGBColor:
    best = rc.ALLOWED_FILLS[0]
    best_d = 1e9
    for c in rc.ALLOWED_FILLS:
        t = _rgb_tuple(c)
        if t is None:
            continue
        d = _dist(rgb, t)
        if d < best_d:
            best_d = d
            best = c
    return best


def _emu_to_in(emu: int) -> float:
    return float(emu) / 914400.0


def _get_fill_rgb(shape) -> RGBColor | None:
    try:
        fill = shape.fill
        if fill.type is None:
            return None
        # solid fill
        if hasattr(fill, "fore_color") and fill.fore_color and fill.fore_color.rgb:
            return fill.fore_color.rgb
    except Exception:  # noqa: BLE001
        return None
    return None


def _set_fill(shape, color: RGBColor) -> bool:
    try:
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = color
        return True
    except Exception:  # noqa: BLE001
        return False


def _guess_role(shape, font_pt: float | None) -> str:
    top = _emu_to_in(shape.top)
    height = _emu_to_in(shape.height)
    if top >= rc.FOOTER_TOP_MIN_IN or (font_pt is not None and font_pt <= rc.LABEL_SIZE + 0.5):
        return "label"
    if (
        height <= rc.MSG_BAR_HEIGHT_MAX_IN
        and top < 2.5
        and font_pt is not None
        and 16 <= font_pt <= 24
    ):
        return "msg_bar"
    if top <= rc.TITLE_TOP_MAX_IN and (
        font_pt is None or font_pt >= rc.TITLE_SIZE_MIN_PT
    ):
        return "title"
    return "body"


def _target_size(role: str) -> int:
    return {
        "title": rc.TITLE_SIZE,
        "msg_bar": rc.MSG_BAR_SIZE,
        "label": rc.LABEL_SIZE,
        "body": rc.BODY_SIZE,
    }.get(role, rc.BODY_SIZE)


def _snap_size(pt: float, role: str) -> int:
    target = _target_size(role)
    if abs(pt - target) <= 1.5:
        return int(round(pt))
    # Prefer role target if far from allowed set
    if int(round(pt)) in rc.ALLOWED_SIZES_PT and abs(pt - target) > 6:
        # Keep sizes that are already on-brand unless wildly off role
        if role == "title" and pt < 24:
            return target
        if role == "msg_bar" and abs(pt - rc.MSG_BAR_SIZE) > 4:
            return target
        return int(round(pt))
    return target


def _iter_runs(shape):
    if not getattr(shape, "has_text_frame", False):
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            yield run


def _first_font_size_pt(shape) -> float | None:
    for run in _iter_runs(shape):
        try:
            if run.font.size is not None:
                return run.font.size.pt
        except Exception:  # noqa: BLE001
            continue
    return None


def _disable_autofit(shape, log: Changelog, slide_idx: int) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    try:
        # Overflow pre-check: if text already long vs box, still disable autofit
        # (ST rule) but record a skip note for potential overflow.
        text_len = len(shape.text_frame.text or "")
        box_area = max(_emu_to_in(shape.width) * _emu_to_in(shape.height), 0.01)
        density = text_len / box_area
        if density > 180:
            log.skip(slide_idx, shape.shape_id, "high_text_density_after_autofit_off")
        if tf.auto_size != MSO_AUTO_SIZE.NONE:
            tf.auto_size = MSO_AUTO_SIZE.NONE
            log.add(slide_idx, shape.shape_id, "autofit", "disabled autofit")
        # word wrap on
        if hasattr(tf, "word_wrap") and tf.word_wrap is False:
            tf.word_wrap = True
            log.add(slide_idx, shape.shape_id, "autofit", "enabled word_wrap")
    except Exception:  # noqa: BLE001
        log.skip(slide_idx, shape.shape_id, "autofit_unavailable")


def _refresh_text_style(
    shape,
    log: Changelog,
    slide_idx: int,
    fill_rgb: RGBColor | None,
    *,
    conservative: bool = False,
) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    font_pt = _first_font_size_pt(shape)
    role = _guess_role(shape, font_pt)
    # Cover / title slides: keep large display sizes; only normalize font + contrast.
    if conservative and font_pt is not None and font_pt >= 24:
        target_pt = int(round(font_pt))
        snap_size = False
    else:
        target_pt = _snap_size(font_pt or float(_target_size(role)), role)
        snap_size = True
    contrast = rc.text_on(fill_rgb) if fill_rgb is not None else None

    for run in _iter_runs(shape):
        # Never touch run.text
        try:
            name = run.font.name
            if name and name not in (rc.FONT, "Arial Narrow", "arial", "ARIAL"):
                run.font.name = rc.FONT
                log.add(slide_idx, shape.shape_id, "font", f"{name} → {rc.FONT}")
            elif not name:
                run.font.name = rc.FONT
                log.add(slide_idx, shape.shape_id, "font", f"set {rc.FONT}")
        except Exception:  # noqa: BLE001
            pass

        if snap_size:
            try:
                cur = run.font.size.pt if run.font.size else None
                if cur is None or abs(cur - target_pt) >= 1.0:
                    run.font.size = Pt(target_pt)
                    log.add(
                        slide_idx,
                        shape.shape_id,
                        "font_size",
                        f"{cur or '?'}pt → {target_pt}pt ({role})",
                    )
            except Exception:  # noqa: BLE001
                pass

        if contrast is not None:
            try:
                cur_c = (
                    run.font.color.rgb
                    if run.font.color and run.font.color.type is not None
                    else None
                )
                if cur_c != contrast:
                    run.font.color.rgb = contrast
                    log.add(slide_idx, shape.shape_id, "contrast", f"text → {contrast}")
            except Exception:  # noqa: BLE001
                pass

        if role == "msg_bar" and not conservative:
            try:
                if not run.font.bold:
                    run.font.bold = True
                    log.add(slide_idx, shape.shape_id, "font_weight", "msg_bar bold")
            except Exception:  # noqa: BLE001
                pass


def _refresh_fill(
    shape, log: Changelog, slide_idx: int, *, conservative: bool = False
) -> RGBColor | None:
    fill_rgb = _get_fill_rgb(shape)
    if fill_rgb is None:
        return None
    # Cover slides: only snap fills that are clearly off-brand (far from palette).
    min_dist = rc.COLOR_SNAP_MIN_DIST * (1.6 if conservative else 1.0)
    tup = _rgb_tuple(fill_rgb)
    if tup is None:
        return fill_rgb
    nearest = _nearest_fill(tup)
    nt = _rgb_tuple(nearest)
    if nt is None:
        return fill_rgb
    d = _dist(tup, nt)
    if d >= min_dist and nearest != fill_rgb:
        if _set_fill(shape, nearest):
            log.add(
                slide_idx,
                shape.shape_id,
                "color",
                f"fill {tup} → {nt} (d={d:.0f})",
            )
            return nearest
    return fill_rgb


def _refresh_picture(shape, log: Changelog, slide_idx: int) -> None:
    try:
        w = _emu_to_in(shape.width)
        h = _emu_to_in(shape.height)
        if h <= 0 or w <= 0:
            return
        log.add(slide_idx, shape.shape_id, "picture", "aspect preserved (no replace)")
    except Exception:  # noqa: BLE001
        log.skip(slide_idx, shape.shape_id, "picture_lock_failed")


def _align_group(shapes: list, axis: str, log: Changelog, slide_idx: int) -> None:
    """Snap left/top of similar shapes when within 2x tolerance cluster."""
    if len(shapes) < 2:
        return
    if axis == "left":
        getter = lambda s: s.left  # noqa: E731
        setter = lambda s, v: setattr(s, "left", v)  # noqa: E731
    elif axis == "top":
        getter = lambda s: s.top  # noqa: E731
        setter = lambda s, v: setattr(s, "top", v)  # noqa: E731
    else:
        return

    tol = int(Inches(rc.ALIGN_TOLERANCE_IN))
    vals_sorted = sorted(getter(s) for s in shapes)
    median = vals_sorted[len(vals_sorted) // 2]
    cluster = [s for s in shapes if abs(getter(s) - median) <= tol * 2]
    if len(cluster) < 2:
        return
    target = int(sum(getter(s) for s in cluster) / len(cluster))
    for s in cluster:
        if abs(getter(s) - target) > tol // 2 and abs(getter(s) - target) <= tol * 2:
            old = getter(s)
            setter(s, target)
            log.add(
                slide_idx,
                s.shape_id,
                "align",
                f"{axis} {_emu_to_in(old):.3f}in → {_emu_to_in(target):.3f}in",
            )


def _equalize_column_widths(shapes: list, log: Changelog, slide_idx: int) -> None:
    if len(shapes) < 2:
        return
    tol = int(Inches(rc.ALIGN_TOLERANCE_IN))
    by_row: dict[int, list] = {}
    for s in shapes:
        key = int(round(s.top / max(tol, 1)))
        by_row.setdefault(key, []).append(s)
    for row in by_row.values():
        if len(row) < 2:
            continue
        widths = [s.width for s in row]
        avg = sum(widths) / len(widths)
        if avg <= 0:
            continue
        if max(abs(w - avg) / avg for w in widths) > 0.15:
            continue
        target = int(avg)
        for s in row:
            if abs(s.width - target) > int(Inches(0.02)):
                old = s.width
                s.width = target
                log.add(
                    slide_idx,
                    s.shape_id,
                    "symmetry",
                    f"width {_emu_to_in(old):.3f} → {_emu_to_in(target):.3f}in",
                )


def _slide_looks_like_cover(slide, slide_idx: int) -> bool:
    """Heuristic: first slide or large dark full-bleed field ≈ presentation title."""
    if slide_idx == 1:
        return True
    try:
        for shape in slide.shapes:
            try:
                w = _emu_to_in(shape.width)
                h = _emu_to_in(shape.height)
                if w >= 12.5 and h >= 6.5:
                    fill = _get_fill_rgb(shape)
                    tup = _rgb_tuple(fill) if fill else None
                    if tup and (tup[0] + tup[1] + tup[2]) < 180:
                        return True
            except Exception:  # noqa: BLE001
                continue
    except Exception:  # noqa: BLE001
        pass
    return False


def _process_shape(
    shape, log: Changelog, slide_idx: int, *, conservative: bool = False
) -> None:
    try:
        stype = shape.shape_type
    except Exception:  # noqa: BLE001
        log.skip(slide_idx, getattr(shape, "shape_id", None), "shape_type_unreadable")
        return

    if stype == MSO_SHAPE_TYPE.GROUP:
        try:
            for child in shape.shapes:
                _process_shape(child, log, slide_idx, conservative=conservative)
        except Exception:  # noqa: BLE001
            log.skip(slide_idx, shape.shape_id, "group_unreadable")
        return

    if stype == MSO_SHAPE_TYPE.CHART:
        log.skip(slide_idx, shape.shape_id, "chart_skipped")
        return
    if stype == MSO_SHAPE_TYPE.TABLE:
        # Style table fonts lightly without changing cell text
        try:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name and run.font.name not in (rc.FONT, "Arial Narrow"):
                                old = run.font.name
                                run.font.name = rc.FONT
                                log.add(slide_idx, shape.shape_id, "font", f"table {old} → {rc.FONT}")
        except Exception:  # noqa: BLE001
            log.skip(slide_idx, shape.shape_id, "table_style_failed")
        return
    if stype == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
        log.skip(slide_idx, shape.shape_id, "ole_skipped")
        return

    if stype == MSO_SHAPE_TYPE.PICTURE:
        _refresh_picture(shape, log, slide_idx)
        return

    fill_rgb = _refresh_fill(shape, log, slide_idx, conservative=conservative)
    _disable_autofit(shape, log, slide_idx)
    _refresh_text_style(
        shape, log, slide_idx, fill_rgb, conservative=conservative
    )

    # Light line cleanup for rectangles with thick borders (skip on cover)
    if conservative:
        return
    try:
        if shape.line and shape.line.width and shape.line.width.pt and shape.line.width.pt > 2.5:
            old = shape.line.width.pt
            shape.line.width = Pt(1.0)
            log.add(slide_idx, shape.shape_id, "line", f"width {old:.1f} → 1.0pt")
    except Exception:  # noqa: BLE001
        pass


def refresh_presentation(src: Path, dest: Path, changelog_path: Path | None = None) -> Changelog:
    """Apply standard-tier visual refresh. Preserves page count and all text."""
    log = Changelog(source=src.name)
    prs = Presentation(str(src))
    log.slide_count = len(prs.slides)

    # Aspect warning
    try:
        w_in = _emu_to_in(prs.slide_width)
        h_in = _emu_to_in(prs.slide_height)
        ratio = w_in / h_in if h_in else 0
        if abs(ratio - (16 / 9)) > 0.05:
            log.warnings.append(f"non_16_9_aspect:{ratio:.3f}")
    except Exception:  # noqa: BLE001
        pass

    for si, slide in enumerate(prs.slides, start=1):
        cover = _slide_looks_like_cover(slide, si)
        if cover:
            log.warnings.append(f"slide_{si}_cover_conservative")
        auto_shapes = []
        for shape in slide.shapes:
            _process_shape(shape, log, si, conservative=cover)
            try:
                if shape.shape_type in (
                    MSO_SHAPE_TYPE.AUTO_SHAPE,
                    MSO_SHAPE_TYPE.TEXT_BOX,
                    MSO_SHAPE_TYPE.PLACEHOLDER,
                ):
                    auto_shapes.append(shape)
            except Exception:  # noqa: BLE001
                continue
        # Alignment / symmetry — skip on cover title slides
        if not cover:
            _align_group(auto_shapes, "left", log, si)
            _align_group(auto_shapes, "top", log, si)
            _equalize_column_widths(auto_shapes, log, si)

    dest.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dest))
    log.finalize()
    out = changelog_path or (dest.parent / "refresh_changelog.json")
    out.write_text(json.dumps(log.to_dict(), ensure_ascii=False, indent=2), encoding="utf-8")
    return log
