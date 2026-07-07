"""ST brand helper library for python-pptx.

Encodes the rules from skills/st-ppt-brand (palette, typography, message bar,
Title-Only content slides, card rows, and diagram primitives) so generated decks
are on-brand by construction. Import this from the deck-build script.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---- Official ST palette (Green Vogue / Gold / Picton Blue / White) ----
ST_DARK_BLUE  = RGBColor(0x03, 0x23, 0x4B)   # Green Vogue
ST_YELLOW     = RGBColor(0xFF, 0xD2, 0x00)   # Gold
ST_LIGHT_BLUE = RGBColor(0x3C, 0xB4, 0xE6)   # Picton Blue
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_1        = RGBColor(0xEE, 0xEF, 0xF1)
GRAY_2        = RGBColor(0xDB, 0xDE, 0xE1)
GRAY_3        = RGBColor(0xC0, 0xC8, 0xD2)
SLATE         = RGBColor(0x42, 0x59, 0x78)   # first shade of dark blue

# Dark-blue shading ramp (graded headers / process steps), darkest first
RAMP = [ST_DARK_BLUE, SLATE, RGBColor(0x80, 0x91, 0xA5), RGBColor(0xC0, 0xC9, 0xCE)]

FONT = "Arial"
FONT_DIAGRAM = "Arial Narrow"  # block-diagram labels only (brand spec)
SLIDE_W = 13.333
SLIDE_H = 7.5

# Typography sampled from Presentation template.potx (slideMaster txStyles + layouts).
# sz values in OOXML are hundredths of a point (3600 = 36 pt).
TITLE_SIZE = 36           # content slide title (titleStyle lvl1, regular)
TITLE_BOLD = False
MSG_BAR_SIZE = 20         # key message bar (bodyStyle lvl2)
BODY_SIZE = 14            # box / bullet body (template content examples)
BODY_L1_SIZE = 24         # first-level emphasis / punchline
BODY_L2_SIZE = 20
BODY_L3_SIZE = 18         # subtitles, secondary headings
BODY_L4_SIZE = 16
CAPTION_SIZE = 13         # dense card / row text (template slide examples)
HEADING_SIZE = 17         # card / column headers in multi-box layouts
SUBTITLE_SIZE = 18
LABEL_SIZE = 11           # footer, slide number, axis labels
FOOTNOTE_SIZE = 8         # closing legal text
AGENDA_TOPIC_SIZE = 28    # Agenda layout numbered tiles + topics
PRESENTATION_TITLE_SIZE = 36
SECTION_TITLE_SIZE = 36
CLOSING_TAGLINE_SIZE = 32
SUB_SIZE = 12             # secondary line under a box header


# Dark fills: white text. Light fills: ST Dark Blue text.
# RAMP[2] (#8091A5) is mid-tone — white text for readable contrast.
_DARK_FILLS = frozenset({ST_DARK_BLUE, SLATE, RAMP[2]})


def ramp_text(step):
    return text_on(RAMP[min(step, 3)])


def text_on(fill_color):
    """Mandatory contrast: white on dark fills; ST Dark Blue on light fills."""
    if fill_color in _DARK_FILLS:
        return WHITE
    return ST_DARK_BLUE


def new_deck():
    """16:9 presentation per ST template geometry."""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def title_only_slide(prs):
    """Reproduce the 'Title Only' layout: blank slide; place shapes yourself."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _place_logo(slide, logo_path, left, top, height=0.48):
    if logo_path:
        slide.shapes.add_picture(logo_path, Inches(left), Inches(top), height=Inches(height))


def presentation_title_slide(prs, title, presenter=None, logo_path=None):
    """Main / presentation title — navy field, left yellow accent, white title (see special-slides.md)."""
    slide = title_only_slide(prs)
    _slide_bg(slide, ST_DARK_BLUE)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                    Inches(0.14), Inches(SLIDE_H))
    fill(accent, ST_YELLOW)
    tb = slide.shapes.add_textbox(Inches(1.1), Inches(2.55), Inches(9.5), Inches(1.15))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = title
    _style(tf, WHITE, PRESENTATION_TITLE_SIZE, bold=True, align=PP_ALIGN.LEFT)
    if presenter:
        ptb = slide.shapes.add_textbox(Inches(1.1), Inches(3.85), Inches(8.0), Inches(0.55))
        ptf = ptb.text_frame
        no_autofit(ptf)
        ptf.text = presenter
        _style(ptf, WHITE, SUBTITLE_SIZE, bold=False, align=PP_ALIGN.LEFT)
    _place_logo(slide, logo_path, 11.35, 0.38, height=0.5)
    return slide


def agenda_slide(prs, topics, title="Agenda", logo_path=None, columns=2):
    """Agenda / table of contents — white field, yellow numbered tiles (see special-slides.md)."""
    slide = title_only_slide(prs)
    _slide_bg(slide, WHITE)
    tb = slide.shapes.add_textbox(Inches(9.6), Inches(0.32), Inches(3.4), Inches(0.9))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = title
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    _style(tf, ST_DARK_BLUE, TITLE_SIZE, bold=TITLE_BOLD, align=PP_ALIGN.RIGHT)
    n = len(topics)
    cols = max(1, min(columns, 2))
    rows = (n + cols - 1) // cols
    col_x = [1.35, 7.05]
    row_y0 = 1.55
    row_gap = 1.05
    tile = 0.42
    for i, topic in enumerate(topics):
        col = i // rows if rows else 0
        row = i % rows if rows else 0
        if col >= cols:
            col = cols - 1
        x = col_x[col]
        y = row_y0 + row * row_gap
        sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                    Inches(tile), Inches(tile))
        fill(sq, ST_YELLOW)
        stf = sq.text_frame
        no_autofit(stf)
        stf.vertical_anchor = MSO_ANCHOR.MIDDLE
        stf.text = str(i + 1)
        for p in stf.paragraphs:
            p.alignment = PP_ALIGN.CENTER
        _style(stf, ST_DARK_BLUE, AGENDA_TOPIC_SIZE, bold=True, align=PP_ALIGN.CENTER)
        label = slide.shapes.add_textbox(Inches(x + tile + 0.18), Inches(y + 0.06),
                                         Inches(4.8), Inches(0.38))
        ltf = label.text_frame
        no_autofit(ltf)
        ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ltf.text = topic
        _style(ltf, ST_DARK_BLUE, AGENDA_TOPIC_SIZE, bold=False, align=PP_ALIGN.LEFT)
    _place_logo(slide, logo_path, 0.45, 6.55, height=0.42)
    return slide


def section_title_slide(prs, title, image_path=None, logo_path=None):
    """Section break — navy field, top yellow bar, optional hero image (see special-slides.md)."""
    slide = title_only_slide(prs)
    _slide_bg(slide, ST_DARK_BLUE)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.35), Inches(0),
                                 Inches(SLIDE_W - 3.35), Inches(1.15))
    fill(bar, ST_YELLOW)
    tf = bar.text_frame
    no_autofit(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.35)
    tf.text = title
    _style(tf, ST_DARK_BLUE, SECTION_TITLE_SIZE, bold=True, align=PP_ALIGN.LEFT)
    if image_path:
        slide.shapes.add_picture(image_path, Inches(2.2), Inches(1.65),
                                 width=Inches(8.9), height=Inches(4.35))
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(6.72),
                                  Inches(12.4), Inches(0.02))
    fill(rule, GRAY_3)
    _place_logo(slide, logo_path, 0.45, 6.85, height=0.38)
    return slide


def _place_image_or_placeholder(slide, image_path, x, y, w, h, label="Image"):
    """Use a real ST-owned image when available; otherwise leave a gray placeholder."""
    if image_path:
        slide.shapes.add_picture(image_path, Inches(x), Inches(y),
                                 width=Inches(w), height=Inches(h))
        return None
    ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    fill(ph, GRAY_2)
    ph.line.color.rgb = GRAY_3
    ph.line.width = Pt(1)
    tf = ph.text_frame
    no_autofit(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.text = label
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
    _style(tf, ST_DARK_BLUE, LABEL_SIZE, bold=False, align=PP_ALIGN.CENTER)
    return ph


def left_image_icon_rows_slide(
    prs,
    title,
    rows,
    punchline=None,
    img_path=None,
    logo_path=None,
    img_placeholder="Image",
):
    """Left hero image + yellow icon tiles + gray statement rows (see layout-library #14)."""
    slide = title_only_slide(prs)
    corner_accent(slide)
    img_w = 4.15
    _place_image_or_placeholder(slide, img_path, 0, 0, img_w, SLIDE_H, label=img_placeholder)
    _place_logo(slide, logo_path, 0.45, 6.55, height=0.4)
    rx = img_w + 0.35
    rw = SLIDE_W - rx - 0.35
    tb = slide.shapes.add_textbox(Inches(rx), Inches(0.32), Inches(rw), Inches(0.7))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = title
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    _style(tf, ST_DARK_BLUE, TITLE_SIZE, bold=TITLE_BOLD, align=PP_ALIGN.LEFT)
    row_h = 0.52
    gap = 0.14
    y0 = 1.25
    tile = 0.48
    for i, row in enumerate(rows):
        y = y0 + i * (row_h + gap)
        sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(rx), Inches(y),
                                    Inches(tile), Inches(row_h))
        fill(sq, ST_YELLOW)
        icon = row.get("icon_path")
        if icon:
            slide.shapes.add_picture(icon, Inches(rx + 0.06), Inches(y + 0.06),
                                     width=Inches(tile - 0.12), height=Inches(row_h - 0.12))
        bar_x = rx + tile + 0.12
        bar_w = rw - tile - 0.12
        box(slide, bar_x, y, bar_w, row_h, row.get("text", ""), GRAY_1,
            ST_DARK_BLUE, size=CAPTION_SIZE, bold=False, align=PP_ALIGN.LEFT)
    if punchline:
        ptb = slide.shapes.add_textbox(Inches(rx), Inches(5.85), Inches(rw), Inches(1.2))
        ptf = ptb.text_frame
        no_autofit(ptf)
        ptf.word_wrap = True
        ptf.text = punchline
        _style(ptf, ST_DARK_BLUE, BODY_L1_SIZE, bold=True, align=PP_ALIGN.LEFT)
    return slide


def left_image_tiered_list_slide(
    prs,
    title,
    message,
    categories,
    img_path=None,
    logo_path=None,
    img_placeholder="Image",
):
    """Left hero + overlapping navy message bar + yellow/gray category rows (layout-library #16)."""
    slide = title_only_slide(prs)
    corner_accent(slide)
    img_w = 5.05
    img_top = 1.35
    img_h = 5.95
    _place_image_or_placeholder(slide, img_path, 0, img_top, img_w, img_h, label=img_placeholder)
    _place_logo(slide, logo_path, 0.45, 6.55, height=0.38)
    rx = img_w + 0.25
    rw = SLIDE_W - rx - 0.35
    tb = slide.shapes.add_textbox(Inches(rx), Inches(0.28), Inches(rw), Inches(0.65))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = title
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    _style(tf, ST_DARK_BLUE, TITLE_SIZE, bold=TITLE_BOLD, align=PP_ALIGN.LEFT)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(0.88),
                                 Inches(10.2), Inches(0.78))
    fill(bar, ST_DARK_BLUE)
    btf = bar.text_frame
    no_autofit(btf)
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    btf.margin_left = Inches(0.25)
    btf.margin_right = Inches(0.25)
    btf.text = message
    _style(btf, WHITE, MSG_BAR_SIZE, bold=True, align=PP_ALIGN.LEFT)
    y = 1.45
    head_w = 1.05
    gap = 0.12
    for cat in categories:
        body_h = max(0.95, 0.28 + 0.22 * len(cat.get("bullets", [])))
        box(slide, rx, y, head_w, body_h, cat.get("title", ""), ST_YELLOW,
            ST_DARK_BLUE, size=CAPTION_SIZE, bold=True, align=PP_ALIGN.LEFT)
        bullet_box(slide, rx + head_w + gap, y, rw - head_w - gap, body_h,
                   cat.get("bullets", []), shade=GRAY_1, size=BODY_SIZE)
        y += body_h + gap
    return slide


def migration_timeline_circles_slide(
    prs,
    title,
    subtitle=None,
    steps=None,
    logo_path=None,
):
    """Wave timeline with alternating circles + yellow callout markers (layout-library #15)."""
    slide = title_only_slide(prs)
    corner_accent(slide)
    steps = steps or []
    tb = slide.shapes.add_textbox(Inches(7.8), Inches(0.28), Inches(5.2), Inches(0.95))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = title
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    _style(tf, ST_DARK_BLUE, TITLE_SIZE, bold=TITLE_BOLD, align=PP_ALIGN.RIGHT)
    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.alignment = PP_ALIGN.RIGHT
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(BODY_L3_SIZE)
            r.font.color.rgb = ST_DARK_BLUE
    n = len(steps)
    if not n:
        return slide
    d = 1.75
    x0 = 0.75
    span = SLIDE_W - 1.5
    gap = (span - d * n) / (n - 1) if n > 1 else 0
    cy = 3.55
    arrow(slide, x0, cy, x0 + span, cy, color=ST_DARK_BLUE, width=1.2, dashed=True)
    for i, step in enumerate(steps):
        x = x0 + i * (d + gap)
        col = ST_DARK_BLUE if i % 2 == 0 else RAMP[2]
        tc = text_on(col)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(cy - d / 2),
                                      Inches(d), Inches(d))
        fill(circ, col)
        ctf = circ.text_frame
        no_autofit(ctf)
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ctf.word_wrap = True
        ctf.text = step.get("label", "")
        for p in ctf.paragraphs:
            p.alignment = PP_ALIGN.CENTER
        _style(ctf, tc, CAPTION_SIZE, bold=True, align=PP_ALIGN.CENTER)
        callout = step.get("callout")
        if callout:
            above = step.get("callout_pos", "above") != "below"
            my = cy - d / 2 - 0.55 if above else cy + d / 2 + 0.18
            dot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + d / 2 - 0.03),
                                         Inches(my - 0.18 if above else cy + d / 2),
                                         Inches(0.06), Inches(0.28 if above else 0.22))
            fill(dot, ST_YELLOW)
            label(slide, x + d / 2 - 1.35, my - 0.15 if above else my,
                  2.7, callout, size=LABEL_SIZE, bold=False, align=PP_ALIGN.CENTER)
    _place_logo(slide, logo_path, 0.45, 6.55, height=0.38)
    return slide


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def no_autofit(tf):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE


def _style(tf, color, size_pt, bold=False, font=FONT, align=None):
    for p in tf.paragraphs:
        if align is not None:
            p.alignment = align
        for r in p.runs:
            r.font.name = font
            r.font.size = Pt(size_pt)
            r.font.bold = bold
            r.font.color.rgb = color


def corner_accent(slide):
    """Thin ST Dark Blue block, top-right (template accent)."""
    a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.55), Inches(0.0),
                               Inches(0.78), Inches(0.34))
    fill(a, ST_DARK_BLUE)
    return a


def add_title(slide, text, subtitle=None, align=PP_ALIGN.RIGHT, size=None, bold=None):
    if size is None:
        size = TITLE_SIZE
    if bold is None:
        bold = TITLE_BOLD
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(0.28), Inches(12.1), Inches(0.95))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = text
    tf.paragraphs[0].alignment = align
    _style(tf, ST_DARK_BLUE, size, bold=bold)
    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.alignment = align
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(SUBTITLE_SIZE)
            r.font.color.rgb = ST_LIGHT_BLUE
    return tb


def add_message_bar(slide, text, fill_color=ST_LIGHT_BLUE):
    """The signature ST element. Geometry fixed by the template; 20pt Arial only.
    Fill must be ST Yellow / Dark Blue / Light Blue / slate. Text is single color.
    """
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.43),
                                 Inches(9.8), Inches(0.84))
    fill(bar, fill_color)
    tf = bar.text_frame
    no_autofit(tf)
    tf.margin_left = Inches(0.3)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.text = text
    txt = text_on(fill_color)
    _style(tf, txt, MSG_BAR_SIZE, bold=True)
    return bar


def box(slide, x, y, w, h, text, fill_color, text_color=None, size=BODY_SIZE, bold=True,
        align=PP_ALIGN.CENTER, sub=None, sub_size=SUB_SIZE):
    """A rectangle with embedded text. text_color defaults to text_on(fill)."""
    if text_color is None:
        text_color = text_on(fill_color)
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    fill(sh, fill_color)
    tf = sh.text_frame
    no_autofit(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.text = text
    p0 = tf.paragraphs[0]
    p0.alignment = align
    for r in p0.runs:
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = text_color
    if sub:
        p = tf.add_paragraph()
        p.text = sub
        p.alignment = align
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(sub_size)
            r.font.color.rgb = text_color
    return sh


def bullet_box(slide, x, y, w, h, bullets, shade=GRAY_1, size=BODY_SIZE, heading=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    fill(sh, shade)
    tf = sh.text_frame
    no_autofit(tf)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.16)
    tf.margin_top = Inches(0.12)
    tf.margin_right = Inches(0.14)
    first = True
    if heading:
        tf.text = heading
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
        for r in tf.paragraphs[0].runs:
            r.font.name = FONT
            r.font.size = Pt(size + (BODY_L4_SIZE - BODY_SIZE))
            r.font.bold = True
            r.font.color.rgb = ST_DARK_BLUE
        first = False
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = "\u2022 " + b
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.color.rgb = ST_DARK_BLUE
    return sh


def label(slide, x, y, w, text, color=ST_DARK_BLUE, size=LABEL_SIZE, bold=True,
          align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.3))
    tf = tb.text_frame
    no_autofit(tf)
    tf.margin_left = Inches(0.02)
    tf.margin_top = Inches(0)
    tf.text = text
    tf.paragraphs[0].alignment = align
    for r in tf.paragraphs[0].runs:
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def arrow(slide, x1, y1, x2, y2, color=ST_DARK_BLUE, width=2.25, dashed=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                                   Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'),
                             {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    if dashed:
        ln.insert(0, ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    return c


def dashed_container(slide, x, y, w, h, color=SLATE):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.fill.background()
    sh.line.color.rgb = color
    sh.line.width = Pt(1.5)
    ln = sh.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    sh.shadow.inherit = False
    return sh


def add_activation_timeline(
    slide,
    checkpoints,
    top_items=None,
    bottom_items=None,
    x0=0.55,
    x1=12.8,
    y_line=3.95,
):
    """Draw a 2-lane activation timeline similar to GTM launch plans.

    checkpoints: [{"x": 2.4, "label": "Mar 5"}, ...]
    top_items / bottom_items:
      [{"x": 2.6, "title": "Milestone", "note": "optional", "w": 2.0, "color": ST_YELLOW}, ...]
    """
    top_items = top_items or []
    bottom_items = bottom_items or []

    # Main axis
    arrow(slide, x0, y_line, x1, y_line, color=ST_DARK_BLUE, width=2.2)

    # Lane labels
    box(slide, 0.0, y_line - 1.55, 0.42, 1.42, "CONTENT", ST_DARK_BLUE, WHITE, size=LABEL_SIZE)
    box(slide, 0.0, y_line + 0.15, 0.42, 1.42, "PROMO", ST_DARK_BLUE, WHITE, size=LABEL_SIZE)

    # Checkpoints on axis
    for cp in checkpoints:
        x = cp["x"]
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - 0.08), Inches(y_line - 0.08),
                                     Inches(0.16), Inches(0.16))
        fill(dot, WHITE)
        dot.line.color.rgb = ST_DARK_BLUE
        dot.line.width = Pt(1.5)
        if cp.get("label"):
            label(slide, x - 0.35, y_line - 0.42, 0.8, cp["label"], color=ST_DARK_BLUE,
                  size=LABEL_SIZE, bold=False, align=PP_ALIGN.CENTER)

    # Activity cards above axis
    for i, item in enumerate(top_items):
        w = item.get("w", 2.0)
        h = 0.5 if not item.get("note") else 0.82
        y = y_line - 1.2 - (0.62 * (i % 2))
        fc = item.get("color", ST_YELLOW)
        tc = text_on(fc)
        box(slide, item["x"] - w / 2, y, w, h, item["title"], fc, tc, size=CAPTION_SIZE,
            align=PP_ALIGN.LEFT, sub=item.get("note"), sub_size=LABEL_SIZE)
        arrow(slide, item["x"], y + h, item["x"], y_line - 0.12, color=GRAY_3, width=1.2)

    # Activity cards below axis
    for i, item in enumerate(bottom_items):
        w = item.get("w", 2.0)
        h = 0.5 if not item.get("note") else 0.82
        y = y_line + 0.26 + (0.62 * (i % 2))
        fc = item.get("color", GRAY_1)
        tc = text_on(fc)
        box(slide, item["x"] - w / 2, y, w, h, item["title"], fc, tc, size=CAPTION_SIZE,
            align=PP_ALIGN.LEFT, sub=item.get("note"), sub_size=LABEL_SIZE)
        arrow(slide, item["x"], y, item["x"], y_line + 0.12, color=GRAY_3, width=1.2)


def timeline_template_slide(prs, title, message, checkpoints, top_items=None, bottom_items=None):
    """One-call template: title + message bar + 2-lane timeline."""
    slide = title_only_slide(prs)
    corner_accent(slide)
    add_title(slide, title, align=PP_ALIGN.RIGHT)
    add_message_bar(slide, message, fill_color=ST_LIGHT_BLUE)
    add_activation_timeline(slide, checkpoints, top_items=top_items, bottom_items=bottom_items)
    return slide


def add_cards_row(slide, cards, top=2.5, bottom=6.95, gap=0.4, left_margin=0.45,
                  header="yellow", img_ratio=2.35):
    """cards = [{"title", "bullets":[...], "img": path|None}].
    header = 'yellow' (yellow bar, dark-blue text) or 'ramp' (graded dark-blue).
    Lays out N equal cards: header bar + optional image banner + gray bullet box.
    """
    from PIL import Image
    n = len(cards)
    total_w = SLIDE_W - 2 * left_margin
    w = (total_w - gap * (n - 1)) / n
    head_h = 0.5
    img_top = top + head_h + 0.08
    img_h = w / img_ratio
    box_top = img_top + img_h + 0.08
    for i, c in enumerate(cards):
        x = left_margin + i * (w + gap)
        if header == "yellow":
            hf, ht = ST_YELLOW, ST_DARK_BLUE
        else:
            hf, ht = RAMP[min(i, 3)], text_on(RAMP[min(i, 3)])
        box(slide, x, top, w, head_h, c["title"], hf, ht, size=BODY_SIZE,
            align=PP_ALIGN.LEFT)
        bt = box_top
        img = c.get("img")
        if img:
            # crop to a clean banner, then place with locked aspect ratio
            im = Image.open(img)
            iw, ih = im.size
            crop_h = min(ih, int(iw / img_ratio))
            cropped = img + ".banner.png"
            im.crop((0, 0, iw, crop_h)).save(cropped)
            slide.shapes.add_picture(cropped, Inches(x), Inches(img_top),
                                     width=Inches(w), height=Inches(img_h))
        else:
            bt = top + head_h
        bullet_box(slide, x, bt, w, top + (bottom - top) - bt,
                   c.get("bullets", []))


def footer(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(7.12), Inches(12.4), Inches(0.32))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = text
    for r in tf.paragraphs[0].runs:
        r.font.name = FONT
        r.font.size = Pt(LABEL_SIZE)
        r.font.color.rgb = ST_DARK_BLUE
    return tb


CLOSING_FOOTER = (
    "© STMicroelectronics - All rights reserved.\n"
    "ST logo is a trademark or a registered trademark of STMicroelectronics "
    "International NV or its affiliates in the EU and/or other countries.\n"
    "For additional information about ST trademarks, please refer to "
    "www.st.com/trademarks.\n"
    "All other product or service names are the property of their respective owners."
)


def closing_slide(prs, logo_path=None, tagline="Our technology starts with You"):
    """Mandatory ST closing / trademark slide for external decks (see compliance checklist)."""
    slide = title_only_slide(prs)
    corner_accent(slide)
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(5.85)
    )
    fill(panel, ST_DARK_BLUE)
    if logo_path:
        slide.shapes.add_picture(logo_path, Inches(0.55), Inches(0.45), height=Inches(0.55))
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.35), Inches(11.7), Inches(1.4))
    tf = tb.text_frame
    no_autofit(tf)
    tf.text = tagline
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _style(tf, WHITE, CLOSING_TAGLINE_SIZE, bold=True)
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.95), Inches(SLIDE_W), Inches(1.55)
    )
    fill(band, ST_YELLOW)
    ftb = slide.shapes.add_textbox(Inches(0.45), Inches(6.05), Inches(12.4), Inches(1.35))
    ftf = ftb.text_frame
    no_autofit(ftf)
    ftf.word_wrap = True
    ftf.text = CLOSING_FOOTER
    for p in ftf.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(2)
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(FOOTNOTE_SIZE)
            r.font.color.rgb = ST_DARK_BLUE
    return slide
